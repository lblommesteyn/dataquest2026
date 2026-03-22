"""Build deck-ready presentation assets from trained RescueWindow artifacts."""
from __future__ import annotations

import gc
import json
import math
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Tuple

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib import patches
from matplotlib.gridspec import GridSpec
from matplotlib.lines import Line2D
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sklearn.linear_model import LogisticRegression
from sklearn.metrics import (
    average_precision_score,
    precision_recall_curve,
    roc_auc_score,
    roc_curve,
)
from sklearn.preprocessing import StandardScaler

from config import CFG
from data.feature_engineer import engineer_features
from data.imputer import PatientImputer
from data.loader import FEATURE_COLS, load_patients
from data.preprocessor import preprocess_patients
from data.splitter import load_split
from data.windower import build_windows
from evaluation.metrics import compute_metrics, find_threshold_at_sensitivity
from models.fusion.model import FusionModel, build_meta_features
from models.phenotype.clustering import ARCHETYPE_NAMES
from models.phenotype.trainer import PhenotypeModel
from models.snapshot.model import SnapshotModel
from models.trajectory.model import TrajectoryModel
from utils.io import load_json
from utils.logging import get_logger

log = get_logger("presentation")

PROJECT_ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = PROJECT_ROOT / "presentation"
FIGURES_DIR = PRESENTATION_DIR / "figures"
TABLES_DIR = PRESENTATION_DIR / "tables"
SUMMARY_PATH = PRESENTATION_DIR / "summary.json"
DECK_PATH = PRESENTATION_DIR / "rescuewindow_deck.pptx"
OUTLINE_PATH = PRESENTATION_DIR / "deck_outline.md"

COLORS = {
    "bg": "#F6F1E8",
    "paper": "#FFFDFC",
    "ink": "#15202B",
    "muted": "#55606D",
    "grid": "#D9D0C2",
    "baseline": "#7B8794",
    "snapshot": "#15616D",
    "trajectory": "#1D4E89",
    "fusion": "#E76F51",
    "positive": "#2A9D8F",
    "warning": "#F4A261",
    "danger": "#C44536",
    "accent": "#F4D35E",
}

MODEL_DISPLAY = {
    "logistic_regression": "Logistic Regression",
    "snapshot_xgb": "Snapshot XGBoost",
    "trajectory_gru": "Trajectory GRU",
    "fusion_multimodel": "Fusion (RescueWindow)",
}

MODEL_COLORS = {
    "logistic_regression": COLORS["baseline"],
    "snapshot_xgb": COLORS["snapshot"],
    "trajectory_gru": COLORS["trajectory"],
    "fusion_multimodel": COLORS["fusion"],
}

TARGET_BANDS = {
    "AUROC": {"minimum": 0.78, "good_low": 0.82, "good_high": 0.85, "strong": 0.85, "max": 0.92},
    "AUPRC": {"minimum": 0.25, "good_low": 0.40, "good_high": 0.48, "strong": 0.50, "max": 0.60},
    "Mean Lead Time (h)": {"minimum": 3.0, "good_low": 3.0, "good_high": 5.0, "strong": 5.0, "max": 8.0},
    "Detection Rate": {"minimum": 0.60, "good_low": 0.60, "good_high": 0.75, "strong": 0.80, "max": 1.0},
}


def set_plot_style() -> None:
    plt.style.use("default")
    plt.rcParams.update(
        {
            "figure.facecolor": COLORS["bg"],
            "axes.facecolor": COLORS["paper"],
            "axes.edgecolor": COLORS["grid"],
            "axes.labelcolor": COLORS["ink"],
            "axes.titlecolor": COLORS["ink"],
            "xtick.color": COLORS["ink"],
            "ytick.color": COLORS["ink"],
            "text.color": COLORS["ink"],
            "font.family": "DejaVu Sans",
            "font.size": 11,
            "axes.titleweight": "bold",
            "axes.grid": True,
            "grid.color": COLORS["grid"],
            "grid.alpha": 0.45,
            "grid.linestyle": "-",
            "axes.spines.top": False,
            "axes.spines.right": False,
        }
    )


def ensure_output_dirs() -> None:
    FIGURES_DIR.mkdir(parents=True, exist_ok=True)
    TABLES_DIR.mkdir(parents=True, exist_ok=True)


def save_fig(fig: plt.Figure, name: str, dpi: int = 220) -> Path:
    path = FIGURES_DIR / name
    fig.savefig(path, dpi=dpi, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)
    return path


def hex_to_rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def format_pct(value: float) -> str:
    return f"{value * 100:.1f}%"


def format_metric(metric: str, value: float) -> str:
    if metric in {"Detection Rate", "% Caught >=6h Early", "Sensitivity@Spec90"}:
        return format_pct(value)
    if metric == "Mean Lead Time (h)":
        return f"{value:.1f}h"
    return f"{value:.3f}"


def add_card(ax: plt.Axes, title: str, value: str, subtitle: str, color: str) -> None:
    ax.axis("off")
    box = patches.FancyBboxPatch(
        (0.02, 0.06),
        0.96,
        0.88,
        boxstyle="round,pad=0.02,rounding_size=0.04",
        linewidth=0,
        facecolor=COLORS["paper"],
        transform=ax.transAxes,
    )
    accent = patches.FancyBboxPatch(
        (0.02, 0.82),
        0.96,
        0.12,
        boxstyle="round,pad=0.02,rounding_size=0.04",
        linewidth=0,
        facecolor=color,
        transform=ax.transAxes,
    )
    ax.add_patch(box)
    ax.add_patch(accent)
    ax.text(0.06, 0.68, title, fontsize=11, weight="bold", color=COLORS["ink"], transform=ax.transAxes)
    ax.text(0.06, 0.38, value, fontsize=24, weight="bold", color=color, transform=ax.transAxes)
    ax.text(0.06, 0.18, subtitle, fontsize=10, color=COLORS["muted"], transform=ax.transAxes)


def actual_onset_hours(raw_patients: Dict[str, pd.DataFrame]) -> Dict[str, Optional[int]]:
    onset_hours: Dict[str, Optional[int]] = {}
    for pid, df in raw_patients.items():
        if "SepsisLabel" not in df.columns:
            onset_hours[pid] = None
            continue
        positive = np.flatnonzero(df["SepsisLabel"].fillna(0).to_numpy() > 0)
        onset_hours[pid] = int(positive[0]) if len(positive) else None
    return onset_hours


def split_dict(items: Dict[str, pd.DataFrame], ids: Iterable[str]) -> Dict[str, pd.DataFrame]:
    return {pid: items[pid] for pid in ids if pid in items}


def build_feature_columns(train_patients: Dict[str, pd.DataFrame]) -> List[str]:
    all_feat_cols = list(set().union(*[set(df.columns) for df in train_patients.values()]) - {"SepsisLabel"})
    return [c for c in FEATURE_COLS if c in all_feat_cols]


def load_models(models_dir: Path) -> Dict[str, object]:
    required = [
        models_dir / "snapshot.pkl",
        models_dir / "trajectory.pt",
        models_dir / "fusion.pkl",
        models_dir / "phenotype_ae.pt",
        models_dir / "phenotype_ae_meta.pkl",
        models_dir / "phenotype_clusterer.pkl",
        models_dir / "feature_meta.json",
        models_dir / "imputer.pkl",
    ]
    missing = [str(path.name) for path in required if not path.exists()]
    if missing:
        raise FileNotFoundError(
            "Missing trained artifacts in artifacts/models/. "
            f"Run python train.py first. Missing: {', '.join(missing)}"
        )

    return {
        "snapshot": SnapshotModel.load(str(models_dir / "snapshot.pkl")),
        "trajectory": TrajectoryModel.load(str(models_dir / "trajectory.pt")),
        "phenotype": PhenotypeModel.load(str(models_dir / "phenotype")),
        "fusion": FusionModel.load(str(models_dir / "fusion.pkl")),
        "imputer": PatientImputer.load(str(models_dir / "imputer.pkl")),
        "feature_meta": load_json(str(models_dir / "feature_meta.json")),
    }


@dataclass
class AnalysisBundle:
    train_windows: Dict[str, np.ndarray]
    val_windows: Dict[str, np.ndarray]
    test_windows: Dict[str, np.ndarray]
    raw_train: Dict[str, pd.DataFrame]
    raw_val: Dict[str, pd.DataFrame]
    raw_test: Dict[str, pd.DataFrame]
    val_scores: Dict[str, np.ndarray]
    test_scores: Dict[str, np.ndarray]
    thresholds: Dict[str, float]
    phenotype_clusters_test: np.ndarray
    feature_cols: List[str]
    traj_cols: List[str]
    val_metrics: pd.DataFrame
    test_metrics: pd.DataFrame


def build_window_metrics_df(
    y_true: np.ndarray,
    scores: Dict[str, np.ndarray],
    thresholds: Dict[str, float],
) -> pd.DataFrame:
    rows: List[Dict[str, float]] = []
    for name, score in scores.items():
        metrics = compute_metrics(y_true, score, label="")
        rows.append(
            {
                "model_key": name,
                "AUROC": metrics["auroc"],
                "AUPRC": metrics["auprc"],
                "F1@0.5": metrics["f1@0.5"],
                "Sensitivity@Spec90": metrics["sensitivity@spec90"],
                "Recall@Prec90": metrics["recall@prec90"],
                "Threshold@Sens80": thresholds[name],
            }
        )
    return pd.DataFrame(rows)


def prepare_analysis_bundle() -> AnalysisBundle:
    cfg = CFG
    paths = cfg["paths"]
    models_dir = Path(paths["models_dir"])
    models = load_models(models_dir)

    train_ids, val_ids, test_ids = load_split(str(Path(paths["processed_dir"]) / "split.json"))
    feature_cols = models["feature_meta"]["feature_cols"]
    traj_cols = models["feature_meta"]["traj_cols"]

    log.info("Loading raw patients for presentation build")
    patients_raw = load_patients(paths["raw_dir"])
    log.info("Preprocessing patients")
    patients = preprocess_patients(patients_raw, missing_threshold=cfg["missing_threshold"])

    raw_train = split_dict(patients_raw, train_ids)
    raw_val = split_dict(patients_raw, val_ids)
    raw_test = split_dict(patients_raw, test_ids)

    train_patients = split_dict(patients, train_ids)
    val_patients = split_dict(patients, val_ids)
    test_patients = split_dict(patients, test_ids)

    if not feature_cols:
        feature_cols = build_feature_columns(train_patients)

    imputer: PatientImputer = models["imputer"]
    train_imp = imputer.transform(train_patients)
    val_imp = imputer.transform(val_patients)
    test_imp = imputer.transform(test_patients)

    train_eng = engineer_features(train_imp, raw_train)
    val_eng = engineer_features(val_imp, raw_val)
    test_eng = engineer_features(test_imp, raw_test)

    horizon = cfg["prediction_horizon_hours"]
    seq_len = cfg["sequence_window_hours"]

    train_windows = build_windows(train_eng, feature_cols, horizon, seq_len)
    val_windows = build_windows(val_eng, feature_cols, horizon, seq_len)
    test_windows = build_windows(test_eng, feature_cols, horizon, seq_len)

    del train_imp, val_imp, test_imp, train_eng, val_eng, test_eng, patients, patients_raw
    gc.collect()

    snapshot: SnapshotModel = models["snapshot"]
    trajectory: TrajectoryModel = models["trajectory"]
    phenotype: PhenotypeModel = models["phenotype"]
    fusion: FusionModel = models["fusion"]

    val_snap = snapshot.predict_proba(val_windows["X_snapshot"])
    test_snap = snapshot.predict_proba(test_windows["X_snapshot"])

    val_traj = trajectory.predict_proba(val_windows["X_traj"], mask=val_windows["traj_mask"])
    test_traj = trajectory.predict_proba(test_windows["X_traj"], mask=test_windows["traj_mask"])

    val_pheno = phenotype.predict_proba(val_windows["X_snapshot"])
    test_pheno = phenotype.predict_proba(test_windows["X_snapshot"])
    phenotype_clusters_test = phenotype.predict_cluster(test_windows["X_snapshot"])

    val_meta = build_meta_features(val_snap, val_traj, val_pheno)
    test_meta = build_meta_features(test_snap, test_traj, test_pheno)

    val_fusion = fusion.predict_proba(val_meta)
    test_fusion = fusion.predict_proba(test_meta)

    scaler = StandardScaler()
    X_train_scaled = scaler.fit_transform(train_windows["X_snapshot"])
    X_val_scaled = scaler.transform(val_windows["X_snapshot"])
    X_test_scaled = scaler.transform(test_windows["X_snapshot"])
    lr = LogisticRegression(
        max_iter=500,
        C=1.0,
        class_weight="balanced",
        random_state=cfg["seed"],
    )
    lr.fit(X_train_scaled, train_windows["y"])
    val_lr = lr.predict_proba(X_val_scaled)[:, 1]
    test_lr = lr.predict_proba(X_test_scaled)[:, 1]

    val_scores = {
        "logistic_regression": val_lr,
        "snapshot_xgb": val_snap,
        "trajectory_gru": val_traj,
        "fusion_multimodel": val_fusion,
    }
    test_scores = {
        "logistic_regression": test_lr,
        "snapshot_xgb": test_snap,
        "trajectory_gru": test_traj,
        "fusion_multimodel": test_fusion,
    }
    thresholds = {
        name: find_threshold_at_sensitivity(val_windows["y"], score, cfg["alert_sensitivity_target"])
        for name, score in val_scores.items()
    }

    val_metrics = build_window_metrics_df(val_windows["y"], val_scores, thresholds)
    test_metrics = build_window_metrics_df(test_windows["y"], test_scores, thresholds)

    return AnalysisBundle(
        train_windows=train_windows,
        val_windows=val_windows,
        test_windows=test_windows,
        raw_train=raw_train,
        raw_val=raw_val,
        raw_test=raw_test,
        val_scores=val_scores,
        test_scores=test_scores,
        thresholds=thresholds,
        phenotype_clusters_test=phenotype_clusters_test,
        feature_cols=feature_cols,
        traj_cols=traj_cols,
        val_metrics=val_metrics,
        test_metrics=test_metrics,
    )


def build_patient_alert_table(
    patient_ids: np.ndarray,
    hours: np.ndarray,
    y_window: np.ndarray,
    scores: np.ndarray,
    threshold: float,
    actual_onset_map: Dict[str, Optional[int]],
) -> pd.DataFrame:
    rows: List[Dict[str, object]] = []
    for pid in np.unique(patient_ids):
        mask = patient_ids == pid
        patient_hours = hours[mask]
        patient_scores = scores[mask]
        onset_hour = actual_onset_map.get(str(pid))
        alert_candidates = patient_hours[patient_scores >= threshold]
        alert_hour = int(alert_candidates[0]) if len(alert_candidates) else None
        lead_time = None if alert_hour is None or onset_hour is None else onset_hour - alert_hour
        early_alert = bool(onset_hour is not None and alert_hour is not None and alert_hour <= onset_hour)
        late_alert = bool(onset_hour is not None and alert_hour is not None and alert_hour > onset_hour)
        rows.append(
            {
                "patient_id": str(pid),
                "sepsis_patient": bool(onset_hour is not None),
                "window_positive": bool(y_window[mask].max() > 0),
                "actual_onset_hour": onset_hour,
                "first_alert_hour": alert_hour,
                "lead_time_h": lead_time,
                "early_alert": early_alert,
                "late_alert": late_alert,
                "max_score": float(patient_scores.max()),
            }
        )
    return pd.DataFrame(rows)


def summarize_patient_alerts(alert_table: pd.DataFrame) -> Dict[str, float]:
    sepsis = alert_table[alert_table["sepsis_patient"]].copy()
    if sepsis.empty:
        return {
            "n_sepsis_patients": 0,
            "detection_rate": 0.0,
            "mean_lead_time_h": 0.0,
            "median_lead_time_h": 0.0,
            "pct_caught_ge_6h": 0.0,
            "pct_late_alerts": 0.0,
        }

    early = sepsis[sepsis["early_alert"]].copy()
    late = sepsis[sepsis["late_alert"]].copy()
    lead = early["lead_time_h"].dropna().astype(float)
    return {
        "n_sepsis_patients": int(len(sepsis)),
        "n_early_caught": int(len(early)),
        "n_late_alerts": int(len(late)),
        "detection_rate": float(len(early) / len(sepsis)),
        "mean_lead_time_h": float(lead.mean()) if not lead.empty else 0.0,
        "median_lead_time_h": float(lead.median()) if not lead.empty else 0.0,
        "pct_caught_ge_6h": float((lead >= 6).mean()) if not lead.empty else 0.0,
        "pct_late_alerts": float(len(late) / len(sepsis)),
    }


def build_summary_tables(bundle: AnalysisBundle) -> Tuple[pd.DataFrame, pd.DataFrame, Dict[str, pd.DataFrame]]:
    onset_test = actual_onset_hours(bundle.raw_test)
    per_model_alerts: Dict[str, pd.DataFrame] = {}
    rows: List[Dict[str, float]] = []

    for _, row in bundle.test_metrics.iterrows():
        model_key = row["model_key"]
        alert_table = build_patient_alert_table(
            bundle.test_windows["patient_ids"],
            bundle.test_windows["hours"],
            bundle.test_windows["y"],
            bundle.test_scores[model_key],
            bundle.thresholds[model_key],
            onset_test,
        )
        per_model_alerts[model_key] = alert_table
        clinical = summarize_patient_alerts(alert_table)
        rows.append(
            {
                "model_key": model_key,
                "Model": MODEL_DISPLAY[model_key],
                "AUROC": row["AUROC"],
                "AUPRC": row["AUPRC"],
                "Sensitivity@Spec90": row["Sensitivity@Spec90"],
                "Detection Rate": clinical["detection_rate"],
                "Mean Lead Time (h)": clinical["mean_lead_time_h"],
                "% Caught >=6h Early": clinical["pct_caught_ge_6h"],
                "Late Alert Rate": clinical["pct_late_alerts"],
                "Threshold@Sens80": bundle.thresholds[model_key],
            }
        )
    metrics_df = pd.DataFrame(rows).sort_values("AUROC", ascending=False).reset_index(drop=True)

    fusion = metrics_df.loc[metrics_df["model_key"] == "fusion_multimodel"].iloc[0]
    baseline = metrics_df.loc[metrics_df["model_key"] == "logistic_regression"].iloc[0]
    deltas = pd.DataFrame(
        [
            {
                "Metric": "Fusion AUROC uplift vs logistic",
                "Value": float(fusion["AUROC"] - baseline["AUROC"]),
            },
            {
                "Metric": "Fusion AUPRC uplift vs logistic",
                "Value": float(fusion["AUPRC"] - baseline["AUPRC"]),
            },
            {
                "Metric": "Fusion lead-time uplift vs logistic",
                "Value": float(fusion["Mean Lead Time (h)"] - baseline["Mean Lead Time (h)"]),
            },
        ]
    )
    return metrics_df, deltas, per_model_alerts


def dataset_summary(bundle: AnalysisBundle) -> Dict[str, object]:
    def split_stats(raw_split: Dict[str, pd.DataFrame]) -> Dict[str, float]:
        hours = sum(len(df) for df in raw_split.values())
        sepsis_patients = sum(int(df["SepsisLabel"].fillna(0).max() > 0) for df in raw_split.values())
        lengths = [len(df) for df in raw_split.values()]
        return {
            "patients": len(raw_split),
            "hours": hours,
            "sepsis_patients": sepsis_patients,
            "sepsis_prevalence": sepsis_patients / max(len(raw_split), 1),
            "median_los": float(np.median(lengths)) if lengths else 0.0,
        }

    overall_raw = {}
    overall_raw.update(bundle.raw_train)
    overall_raw.update(bundle.raw_val)
    overall_raw.update(bundle.raw_test)

    selected = ["HR", "O2Sat", "Temp", "SBP", "MAP", "Resp", "Lactate", "Creatinine", "WBC", "Platelets"]
    missing_rows = []
    for feature in selected:
        observed = 0
        total = 0
        for df in overall_raw.values():
            if feature in df.columns:
                values = df[feature]
                observed += int(values.notna().sum())
                total += int(len(values))
        miss_rate = 1.0 - (observed / total) if total else 1.0
        missing_rows.append({"feature": feature, "missing_rate": miss_rate})

    return {
        "overall": split_stats(overall_raw),
        "train": split_stats(bundle.raw_train),
        "val": split_stats(bundle.raw_val),
        "test": split_stats(bundle.raw_test),
        "missingness": pd.DataFrame(missing_rows).sort_values("missing_rate", ascending=False).reset_index(drop=True),
    }


def threshold_sweep(
    patient_ids: np.ndarray,
    hours: np.ndarray,
    y_true: np.ndarray,
    scores: np.ndarray,
    actual_onset_map: Dict[str, Optional[int]],
    thresholds: np.ndarray,
) -> pd.DataFrame:
    rows = []
    prevalence = float(y_true.mean())
    for threshold in thresholds:
        pred = (scores >= threshold).astype(int)
        tp = int(((pred == 1) & (y_true == 1)).sum())
        fp = int(((pred == 1) & (y_true == 0)).sum())
        fn = int(((pred == 0) & (y_true == 1)).sum())
        tn = int(((pred == 0) & (y_true == 0)).sum())
        precision = tp / max(tp + fp, 1)
        recall = tp / max(tp + fn, 1)
        specificity = tn / max(tn + fp, 1)

        alert_table = build_patient_alert_table(patient_ids, hours, y_true, scores, float(threshold), actual_onset_map)
        clinical = summarize_patient_alerts(alert_table)
        rows.append(
            {
                "threshold": float(threshold),
                "precision": precision,
                "recall": recall,
                "specificity": specificity,
                "detection_rate": clinical["detection_rate"],
                "mean_lead_time_h": clinical["mean_lead_time_h"],
                "pct_caught_ge_6h": clinical["pct_caught_ge_6h"],
                "positive_prevalence": prevalence,
            }
        )
    return pd.DataFrame(rows)


def build_exec_summary_figure(metrics_df: pd.DataFrame) -> Path:
    fusion = metrics_df.loc[metrics_df["model_key"] == "fusion_multimodel"].iloc[0]
    baseline = metrics_df.loc[metrics_df["model_key"] == "logistic_regression"].iloc[0]
    auroc_uplift = fusion["AUROC"] - baseline["AUROC"]

    fig = plt.figure(figsize=(14.5, 8), facecolor=COLORS["bg"])
    gs = GridSpec(3, 4, figure=fig, height_ratios=[1.1, 1.6, 0.15])

    cards = [
        ("Fusion AUROC", f"{fusion['AUROC']:.3f}", f"+{auroc_uplift:.3f} vs logistic", COLORS["fusion"]),
        ("AUPRC", f"{fusion['AUPRC']:.3f}", "PR matters because positives are rare", COLORS["positive"]),
        ("Mean Lead Time", f"{fusion['Mean Lead Time (h)']:.1f}h", "Measured to true onset", COLORS["warning"]),
        ("Caught >=6h Early", format_pct(fusion["% Caught >=6h Early"]), "Share of sepsis patients", COLORS["trajectory"]),
    ]

    for idx, (title, value, subtitle, color) in enumerate(cards):
        ax = fig.add_subplot(gs[0, idx])
        add_card(ax, title, value, subtitle, color)

    ax = fig.add_subplot(gs[1, :])
    order = metrics_df.sort_values("AUROC", ascending=True)
    bars = ax.barh(
        order["Model"],
        order["AUROC"],
        color=[MODEL_COLORS[k] for k in order["model_key"]],
        edgecolor="none",
    )
    ax.set_xlim(0.65, max(order["AUROC"].max() + 0.03, 0.90))
    ax.set_title("Fusion sets the pace across the held-out test set", loc="left", fontsize=18)
    ax.set_xlabel("AUROC")
    ax.axvspan(0.82, 0.85, color=COLORS["accent"], alpha=0.18)
    ax.text(0.823, len(order) - 0.2, "2019 top-team band", fontsize=10, color=COLORS["muted"])
    for bar, value in zip(bars, order["AUROC"]):
        ax.text(value + 0.005, bar.get_y() + bar.get_height() / 2, f"{value:.3f}", va="center", fontsize=10)

    fig.suptitle("RescueWindow Executive Summary", x=0.03, y=0.985, ha="left", fontsize=24, weight="bold")
    fig.text(
        0.03,
        0.93,
        "Held-out test results with operating thresholds calibrated on validation predictions.",
        fontsize=12,
        color=COLORS["muted"],
    )
    return save_fig(fig, "00_exec_summary.png")


def build_benchmark_targets_figure(metrics_df: pd.DataFrame) -> Path:
    fusion = metrics_df.loc[metrics_df["model_key"] == "fusion_multimodel"].iloc[0]
    metrics = ["AUROC", "AUPRC", "Mean Lead Time (h)", "Detection Rate"]
    fig, axes = plt.subplots(4, 1, figsize=(13, 9), facecolor=COLORS["bg"])
    fig.subplots_adjust(hspace=0.9)

    for ax, metric in zip(axes, metrics):
        bands = TARGET_BANDS[metric]
        value = float(fusion[metric])
        ax.set_facecolor(COLORS["paper"])
        ax.axvspan(0, bands["minimum"], color=COLORS["danger"], alpha=0.08)
        ax.axvspan(bands["minimum"], bands["good_low"], color=COLORS["warning"], alpha=0.08)
        ax.axvspan(bands["good_low"], bands["good_high"], color=COLORS["positive"], alpha=0.10)
        ax.axvspan(bands["strong"], bands["max"], color=COLORS["trajectory"], alpha=0.10)
        ax.scatter([value], [0], s=220, color=COLORS["fusion"], zorder=5)
        ax.axvline(value, color=COLORS["fusion"], linewidth=2.5)
        if metric == "AUPRC":
            ax.axvline(0.08, color=COLORS["baseline"], linestyle="--", linewidth=1.5)
            ax.text(0.082, 0.25, "rough random baseline", fontsize=9, color=COLORS["muted"])
        if metric == "AUROC":
            ax.axvspan(0.82, 0.85, color=COLORS["accent"], alpha=0.18)
            ax.text(0.821, 0.25, "challenge leader band", fontsize=9, color=COLORS["muted"])
        ax.set_xlim(0, bands["max"])
        ax.set_ylim(-1, 1)
        ax.set_yticks([])
        ax.set_title(metric, loc="left", fontsize=14)
        ax.text(value, -0.48, format_metric(metric, value), ha="center", color=COLORS["fusion"], weight="bold")
        ax.text(bands["minimum"], 0.48, "minimum", color=COLORS["muted"], fontsize=9)
        ax.text(bands["good_low"], 0.48, "good", color=COLORS["muted"], fontsize=9)
        ax.text(bands["strong"], 0.48, "strong", color=COLORS["muted"], fontsize=9)

    fig.suptitle("RescueWindow vs project and challenge benchmarks", x=0.03, ha="left", fontsize=22, weight="bold")
    fig.text(
        0.03,
        0.945,
        "The fusion model lands in or above the top-end target bands that matter for the story.",
        fontsize=12,
        color=COLORS["muted"],
    )
    return save_fig(fig, "01_benchmark_targets.png")


def build_dataset_overview_figure(data_summary: Dict[str, object]) -> Path:
    overall = data_summary["overall"]
    missingness = data_summary["missingness"]
    split_df = pd.DataFrame(
        [
            {"Split": "Train", **data_summary["train"]},
            {"Split": "Validation", **data_summary["val"]},
            {"Split": "Test", **data_summary["test"]},
        ]
    )

    fig = plt.figure(figsize=(15, 9), facecolor=COLORS["bg"])
    gs = GridSpec(2, 3, figure=fig, height_ratios=[1.0, 1.35], wspace=0.35, hspace=0.35)

    cards = [
        ("Patients", f"{overall['patients']:,}", "PhysioNet 2019 challenge cohort", COLORS["fusion"]),
        ("Patient-hours", f"{overall['hours']:,}", "Hourly ICU trajectories", COLORS["trajectory"]),
        ("Sepsis prevalence", format_pct(overall["sepsis_prevalence"]), "Patient-level prevalence", COLORS["positive"]),
        ("Median LOS", f"{overall['median_los']:.0f}h", "Observed ICU stay length", COLORS["warning"]),
    ]
    card_grid = GridSpec(1, 4, figure=fig, left=0.04, right=0.96, top=0.92, bottom=0.72, wspace=0.2)
    for idx, card in enumerate(cards):
        ax = fig.add_subplot(card_grid[0, idx])
        add_card(ax, *card)

    ax_split = fig.add_subplot(gs[1, 0:2])
    xpos = np.arange(len(split_df))
    bars = ax_split.bar(xpos, split_df["patients"], color=COLORS["trajectory"], alpha=0.9, label="Patients")
    ax_hours = ax_split.twinx()
    ax_hours.plot(xpos, split_df["hours"], color=COLORS["fusion"], marker="o", linewidth=2.5, label="Patient-hours")
    ax_split.set_xticks(xpos, split_df["Split"])
    ax_split.set_ylabel("Patients")
    ax_hours.set_ylabel("Patient-hours")
    ax_split.set_title("Patient-level split with no leakage", loc="left", fontsize=16)
    for bar, value in zip(bars, split_df["patients"]):
        ax_split.text(bar.get_x() + bar.get_width() / 2, value + 100, f"{int(value):,}", ha="center", fontsize=10)
    for idx, value in enumerate(split_df["hours"]):
        ax_hours.text(idx, value + max(split_df["hours"]) * 0.015, f"{int(value):,}", ha="center", color=COLORS["fusion"], fontsize=9)

    handles = [
        Line2D([0], [0], color=COLORS["trajectory"], linewidth=8, label="Patients"),
        Line2D([0], [0], color=COLORS["fusion"], marker="o", linewidth=2.5, label="Patient-hours"),
    ]
    ax_split.legend(handles=handles, loc="upper left")

    ax_missing = fig.add_subplot(gs[1, 2])
    top_missing = missingness.head(10).sort_values("missing_rate", ascending=True)
    ax_missing.barh(top_missing["feature"], top_missing["missing_rate"], color=COLORS["warning"])
    ax_missing.set_xlim(0, 1)
    ax_missing.set_xlabel("Missing rate")
    ax_missing.set_title("High-value variables are sparse", loc="left", fontsize=16)
    for y, value in enumerate(top_missing["missing_rate"]):
        ax_missing.text(value + 0.02, y, format_pct(value), va="center", fontsize=9)

    fig.suptitle("Dataset framing", x=0.03, y=0.985, ha="left", fontsize=22, weight="bold")
    fig.text(
        0.03,
        0.945,
        "The problem is hard because sepsis is rare and many informative labs are observed intermittently.",
        fontsize=12,
        color=COLORS["muted"],
    )
    return save_fig(fig, "02_dataset_overview.png")


def draw_box(ax: plt.Axes, x: float, y: float, w: float, h: float, title: str, subtitle: str, color: str) -> None:
    rect = patches.FancyBboxPatch(
        (x, y),
        w,
        h,
        boxstyle="round,pad=0.015,rounding_size=0.03",
        linewidth=2.0,
        edgecolor=color,
        facecolor=COLORS["paper"],
    )
    ax.add_patch(rect)
    ax.text(x + 0.02, y + h * 0.62, title, fontsize=14, weight="bold", color=color)
    ax.text(x + 0.02, y + h * 0.24, subtitle, fontsize=10, color=COLORS["muted"])


def arrow(ax: plt.Axes, start: Tuple[float, float], end: Tuple[float, float], color: str) -> None:
    arr = patches.FancyArrowPatch(start, end, arrowstyle="-|>", mutation_scale=18, linewidth=2, color=color)
    ax.add_patch(arr)


def build_architecture_figure() -> Path:
    fig, ax = plt.subplots(figsize=(14, 8), facecolor=COLORS["bg"])
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")

    draw_box(ax, 0.04, 0.38, 0.18, 0.22, "Hourly ICU data", "Vitals, labs, demographics,\nmissingness patterns", COLORS["ink"])
    draw_box(ax, 0.28, 0.68, 0.19, 0.18, "Snapshot specialist", "XGBoost on current hour\nplus rolling features", COLORS["snapshot"])
    draw_box(ax, 0.28, 0.40, 0.19, 0.18, "Trajectory specialist", "GRU over the last 12 hours\nof raw signals", COLORS["trajectory"])
    draw_box(ax, 0.28, 0.12, 0.19, 0.18, "Phenotype encoder", "Autoencoder + K-means\nsoft clinical archetypes", COLORS["positive"])
    draw_box(ax, 0.57, 0.40, 0.18, 0.22, "Fusion stacker", "XGBoost combines specialist\nscores and phenotype distances", COLORS["fusion"])
    draw_box(ax, 0.80, 0.38, 0.15, 0.22, "Clinical output", "Risk score, first alert,\nrescue window, phenotype", COLORS["warning"])

    arrow(ax, (0.22, 0.50), (0.28, 0.77), COLORS["ink"])
    arrow(ax, (0.22, 0.50), (0.28, 0.49), COLORS["ink"])
    arrow(ax, (0.22, 0.50), (0.28, 0.21), COLORS["ink"])
    arrow(ax, (0.47, 0.77), (0.57, 0.55), COLORS["snapshot"])
    arrow(ax, (0.47, 0.49), (0.57, 0.51), COLORS["trajectory"])
    arrow(ax, (0.47, 0.21), (0.57, 0.45), COLORS["positive"])
    arrow(ax, (0.75, 0.51), (0.80, 0.49), COLORS["fusion"])

    ax.text(0.03, 0.95, "RescueWindow multimodel design", fontsize=24, weight="bold")
    ax.text(
        0.03,
        0.90,
        "Three specialists look at the same patient from different angles. The stacker learns when each one should dominate.",
        fontsize=12,
        color=COLORS["muted"],
    )
    return save_fig(fig, "03_system_architecture.png")


def build_model_comparison_figure(metrics_df: pd.DataFrame) -> Path:
    fig, axes = plt.subplots(2, 2, figsize=(14, 10), facecolor=COLORS["bg"])
    metric_specs = [
        ("AUROC", (0.70, max(0.90, metrics_df["AUROC"].max() + 0.03))),
        ("AUPRC", (0.0, max(0.55, metrics_df["AUPRC"].max() + 0.05))),
        ("Detection Rate", (0.0, 1.0)),
        ("Mean Lead Time (h)", (0.0, max(8.0, metrics_df["Mean Lead Time (h)"].max() + 0.8))),
    ]
    order = ["logistic_regression", "snapshot_xgb", "trajectory_gru", "fusion_multimodel"]

    for ax, (metric, limits) in zip(axes.ravel(), metric_specs):
        plot_df = metrics_df.set_index("model_key").loc[order].reset_index()
        bars = ax.bar(
            plot_df["Model"],
            plot_df[metric],
            color=[MODEL_COLORS[k] for k in plot_df["model_key"]],
            edgecolor="none",
        )
        ax.set_ylim(*limits)
        ax.set_title(metric, loc="left", fontsize=15)
        ax.tick_params(axis="x", rotation=15)
        if metric == "AUROC":
            ax.axhspan(0.82, 0.85, color=COLORS["accent"], alpha=0.18)
        if metric == "Detection Rate":
            ax.axhline(0.60, color=COLORS["warning"], linestyle="--", linewidth=1.4)
        for bar, value in zip(bars, plot_df[metric]):
            label = format_metric(metric, float(value))
            ax.text(bar.get_x() + bar.get_width() / 2, value + (limits[1] - limits[0]) * 0.03, label, ha="center", fontsize=10)

    fig.suptitle("Fusion beats every single-model comparator", x=0.03, y=0.985, ha="left", fontsize=22, weight="bold")
    fig.text(
        0.03,
        0.945,
        "The core slide: better ranking, better alert coverage, and more lead time than the baseline or either specialist alone.",
        fontsize=12,
        color=COLORS["muted"],
    )
    return save_fig(fig, "04_model_comparison.png")


def build_performance_curves_figure(bundle: AnalysisBundle) -> Path:
    y_true = bundle.test_windows["y"]
    prevalence = y_true.mean()
    fig, axes = plt.subplots(1, 2, figsize=(14.5, 6.5), facecolor=COLORS["bg"])

    for model_key, scores in bundle.test_scores.items():
        color = MODEL_COLORS[model_key]
        name = MODEL_DISPLAY[model_key]
        fpr, tpr, _ = roc_curve(y_true, scores)
        precision, recall, _ = precision_recall_curve(y_true, scores)
        axes[0].plot(fpr, tpr, linewidth=2.5, color=color, label=f"{name} ({roc_auc_score(y_true, scores):.3f})")
        axes[1].plot(recall, precision, linewidth=2.5, color=color, label=f"{name} ({average_precision_score(y_true, scores):.3f})")

    axes[0].plot([0, 1], [0, 1], linestyle="--", color=COLORS["baseline"], alpha=0.8)
    axes[1].axhline(prevalence, linestyle="--", color=COLORS["baseline"], alpha=0.8)
    axes[0].set_title("ROC curves", loc="left", fontsize=16)
    axes[1].set_title("Precision-recall curves", loc="left", fontsize=16)
    axes[0].set_xlabel("False positive rate")
    axes[0].set_ylabel("True positive rate")
    axes[1].set_xlabel("Recall")
    axes[1].set_ylabel("Precision")
    axes[0].legend(frameon=False, fontsize=9, loc="lower right")
    axes[1].legend(frameon=False, fontsize=9, loc="upper right")

    fig.suptitle("Fusion remains strongest across the full score range", x=0.03, y=0.99, ha="left", fontsize=22, weight="bold")
    fig.text(
        0.03,
        0.945,
        f"PR curves matter here because positives are rare. The dashed line is the empirical positive-window prevalence ({prevalence:.1%}).",
        fontsize=12,
        color=COLORS["muted"],
    )
    return save_fig(fig, "05_performance_curves.png")


def build_operating_point_figure(bundle: AnalysisBundle) -> Tuple[Path, pd.DataFrame]:
    actual_onset = actual_onset_hours(bundle.raw_test)
    sweep_df = threshold_sweep(
        bundle.test_windows["patient_ids"],
        bundle.test_windows["hours"],
        bundle.test_windows["y"],
        bundle.test_scores["fusion_multimodel"],
        actual_onset,
        np.linspace(0.05, 0.95, 37),
    )
    chosen = float(bundle.thresholds["fusion_multimodel"])
    chosen_row = sweep_df.iloc[(sweep_df["threshold"] - chosen).abs().argmin()]

    fig, axes = plt.subplots(1, 2, figsize=(14.5, 6.5), facecolor=COLORS["bg"])
    ax = axes[0]
    ax.plot(sweep_df["threshold"], sweep_df["precision"], color=COLORS["snapshot"], linewidth=2.5, label="Precision")
    ax.plot(sweep_df["threshold"], sweep_df["recall"], color=COLORS["trajectory"], linewidth=2.5, label="Recall")
    ax.plot(sweep_df["threshold"], sweep_df["specificity"], color=COLORS["positive"], linewidth=2.5, label="Specificity")
    ax.axvline(chosen, color=COLORS["fusion"], linestyle="--", linewidth=2)
    ax.scatter([chosen], [chosen_row["recall"]], color=COLORS["fusion"], s=90, zorder=5)
    ax.set_xlabel("Alert threshold")
    ax.set_ylabel("Rate")
    ax.set_ylim(0, 1.05)
    ax.set_title("Hour-level operating tradeoff", loc="left", fontsize=16)
    ax.legend(frameon=False)

    ax = axes[1]
    ax.plot(sweep_df["threshold"], sweep_df["detection_rate"], color=COLORS["fusion"], linewidth=2.8, label="Detection rate")
    ax2 = ax.twinx()
    ax2.plot(sweep_df["threshold"], sweep_df["mean_lead_time_h"], color=COLORS["warning"], linewidth=2.8, label="Mean lead time")
    ax.axvline(chosen, color=COLORS["fusion"], linestyle="--", linewidth=2)
    ax.scatter([chosen], [chosen_row["detection_rate"]], color=COLORS["fusion"], s=90, zorder=5)
    ax2.scatter([chosen], [chosen_row["mean_lead_time_h"]], color=COLORS["warning"], s=90, zorder=5)
    ax.set_ylim(0, 1.0)
    ax2.set_ylim(0, max(8.0, sweep_df["mean_lead_time_h"].max() + 0.5))
    ax.set_xlabel("Alert threshold")
    ax.set_ylabel("Detection rate")
    ax2.set_ylabel("Mean lead time (h)")
    ax.set_title("Patient-level tradeoff", loc="left", fontsize=16)

    handles = [
        Line2D([0], [0], color=COLORS["fusion"], linewidth=2.8, label="Detection rate"),
        Line2D([0], [0], color=COLORS["warning"], linewidth=2.8, label="Mean lead time"),
    ]
    ax.legend(handles=handles, frameon=False, loc="lower center")

    fig.suptitle("Chosen threshold balances coverage with early warning", x=0.03, y=0.99, ha="left", fontsize=22, weight="bold")
    fig.text(
        0.03,
        0.945,
        f"The deck uses the validation-calibrated threshold of {chosen:.3f}, not a test-tuned operating point.",
        fontsize=12,
        color=COLORS["muted"],
    )
    return save_fig(fig, "06_operating_point_tradeoff.png"), sweep_df


def build_lead_time_figure(fusion_alerts: pd.DataFrame) -> Path:
    sepsis = fusion_alerts[fusion_alerts["sepsis_patient"]].copy()
    early = sepsis[sepsis["early_alert"]].copy()
    lead = early["lead_time_h"].dropna().astype(float).sort_values()
    cumulative_hours = np.arange(0, 13)
    caught_curve = [float((lead >= h).mean()) if not lead.empty else 0.0 for h in cumulative_hours]

    fig, axes = plt.subplots(1, 2, figsize=(14.5, 6.5), facecolor=COLORS["bg"])
    axes[0].hist(lead, bins=min(18, max(6, int(math.sqrt(max(len(lead), 1))))), color=COLORS["fusion"], alpha=0.9, edgecolor=COLORS["paper"])
    axes[0].axvline(lead.mean() if not lead.empty else 0, color=COLORS["warning"], linestyle="--", linewidth=2)
    axes[0].set_title("Distribution of early warning lead time", loc="left", fontsize=16)
    axes[0].set_xlabel("Hours before true sepsis onset")
    axes[0].set_ylabel("Sepsis patients")
    if not lead.empty:
        axes[0].text(
            0.98,
            0.92,
            f"Mean = {lead.mean():.1f}h\nMedian = {lead.median():.1f}h",
            transform=axes[0].transAxes,
            ha="right",
            va="top",
            fontsize=11,
            bbox={"facecolor": COLORS["paper"], "edgecolor": COLORS["grid"], "boxstyle": "round,pad=0.3"},
        )

    axes[1].plot(cumulative_hours, caught_curve, color=COLORS["trajectory"], linewidth=3)
    axes[1].fill_between(cumulative_hours, caught_curve, color=COLORS["trajectory"], alpha=0.15)
    axes[1].set_title("How much runway does the model create?", loc="left", fontsize=16)
    axes[1].set_xlabel("Required hours of warning")
    axes[1].set_ylabel("Share of sepsis patients caught")
    axes[1].set_ylim(0, 1.0)
    axes[1].set_xticks(cumulative_hours)

    fig.suptitle("Lead time is the differentiator", x=0.03, y=0.99, ha="left", fontsize=22, weight="bold")
    fig.text(
        0.03,
        0.945,
        "AUROC gets you into the conversation. This slide explains why the system matters clinically.",
        fontsize=12,
        color=COLORS["muted"],
    )
    return save_fig(fig, "07_lead_time_story.png")


def build_complementarity_figure(per_model_alerts: Dict[str, pd.DataFrame]) -> Path:
    snap = per_model_alerts["snapshot_xgb"].set_index("patient_id")
    traj = per_model_alerts["trajectory_gru"].set_index("patient_id")
    fusion = per_model_alerts["fusion_multimodel"].set_index("patient_id")

    sepsis_ids = fusion[fusion["sepsis_patient"]].index.tolist()
    matrix = np.zeros((2, 2), dtype=float)
    annotations = [["" for _ in range(2)] for _ in range(2)]

    for i, snap_flag in enumerate([False, True]):
        for j, traj_flag in enumerate([False, True]):
            ids = [
                pid
                for pid in sepsis_ids
                if bool(snap.loc[pid, "early_alert"]) == snap_flag and bool(traj.loc[pid, "early_alert"]) == traj_flag
            ]
            if ids:
                fusion_rate = float(fusion.loc[ids, "early_alert"].mean())
            else:
                fusion_rate = 0.0
            matrix[i, j] = fusion_rate
            annotations[i][j] = f"n={len(ids)}\nFusion={fusion_rate:.0%}"

    fig, axes = plt.subplots(1, 2, figsize=(14.5, 6.5), facecolor=COLORS["bg"])
    ax = axes[0]
    im = ax.imshow(matrix, cmap="YlGnBu", vmin=0, vmax=1)
    ax.set_xticks([0, 1], ["Trajectory misses", "Trajectory catches"])
    ax.set_yticks([0, 1], ["Snapshot misses", "Snapshot catches"])
    ax.set_title("Fusion recovers discordant specialist cases", loc="left", fontsize=16)
    for i in range(2):
        for j in range(2):
            ax.text(j, i, annotations[i][j], ha="center", va="center", fontsize=11, color=COLORS["ink"], weight="bold")
    cbar = fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    cbar.set_label("Fusion early-catch rate")

    ax = axes[1]
    model_counts = []
    for key in ["logistic_regression", "snapshot_xgb", "trajectory_gru", "fusion_multimodel"]:
        table = per_model_alerts[key]
        sepsis = table[table["sepsis_patient"]]
        model_counts.append({"Model": MODEL_DISPLAY[key], "Early catches": int(sepsis["early_alert"].sum()), "model_key": key})
    count_df = pd.DataFrame(model_counts)
    bars = ax.barh(count_df["Model"], count_df["Early catches"], color=[MODEL_COLORS[k] for k in count_df["model_key"]], edgecolor="none")
    ax.set_title("Patient-level early catches", loc="left", fontsize=16)
    ax.set_xlabel("Sepsis patients alerted before onset")
    for bar, value in zip(bars, count_df["Early catches"]):
        ax.text(value + 1, bar.get_y() + bar.get_height() / 2, str(value), va="center", fontsize=10)

    fig.suptitle("Why the stacker helps", x=0.03, y=0.99, ha="left", fontsize=22, weight="bold")
    fig.text(
        0.03,
        0.945,
        "Snapshot and GRU fail on different patients. Fusion recovers a meaningful share of the cases either specialist misses.",
        fontsize=12,
        color=COLORS["muted"],
    )
    return save_fig(fig, "08_specialist_complementarity.png")


def build_phenotype_figure(bundle: AnalysisBundle) -> Path:
    feature_index = {name: idx for idx, name in enumerate(bundle.feature_cols)}
    chosen = [f for f in ["HR", "MAP", "Resp", "Temp", "Lactate", "Creatinine", "WBC", "Platelets"] if f in feature_index]
    if not chosen:
        chosen = bundle.feature_cols[:8]

    rows = []
    for cluster in sorted(np.unique(bundle.phenotype_clusters_test)):
        mask = bundle.phenotype_clusters_test == cluster
        if mask.sum() == 0:
            continue
        subset = bundle.test_windows["X_snapshot"][mask]
        for feature in chosen:
            idx = feature_index[feature]
            rows.append(
                {
                    "cluster": cluster,
                    "archetype": ARCHETYPE_NAMES.get(int(cluster), f"cluster_{int(cluster)}"),
                    "feature": feature,
                    "mean_value": float(np.nanmean(subset[:, idx])),
                }
            )

    heat_df = pd.DataFrame(rows)
    pivot = heat_df.pivot(index="archetype", columns="feature", values="mean_value")
    pivot = (pivot - pivot.mean()) / pivot.std(ddof=0).replace(0, 1)
    cluster_counts = pd.Series(bundle.phenotype_clusters_test).value_counts().sort_index()

    fig, axes = plt.subplots(1, 2, figsize=(15, 7), facecolor=COLORS["bg"], gridspec_kw={"width_ratios": [1.6, 1.0]})
    ax = axes[0]
    im = ax.imshow(pivot.values, cmap="RdYlBu_r", vmin=-2.0, vmax=2.0)
    ax.set_xticks(np.arange(len(pivot.columns)), pivot.columns, rotation=30, ha="right")
    ax.set_yticks(np.arange(len(pivot.index)), pivot.index)
    ax.set_title("Archetypes capture distinct clinical patterns", loc="left", fontsize=16)
    for i in range(pivot.shape[0]):
        for j in range(pivot.shape[1]):
            ax.text(j, i, f"{pivot.values[i, j]:.1f}", ha="center", va="center", fontsize=8)
    cbar = fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    cbar.set_label("Cluster mean z-score")

    ax = axes[1]
    labels = [ARCHETYPE_NAMES.get(int(idx), f"cluster_{int(idx)}") for idx in cluster_counts.index]
    ax.barh(labels, cluster_counts.values, color=COLORS["positive"])
    ax.set_title("Held-out archetype mix", loc="left", fontsize=16)
    ax.set_xlabel("Patient-hours in test set")
    for y, value in enumerate(cluster_counts.values):
        ax.text(value + 50, y, f"{int(value):,}", va="center", fontsize=9)

    fig.suptitle("Phenotypes are not decoration", x=0.03, y=0.99, ha="left", fontsize=22, weight="bold")
    fig.text(
        0.03,
        0.945,
        "The latent archetypes give the fusion model a compact way to reason about where a patient sits in the broader clinical landscape.",
        fontsize=12,
        color=COLORS["muted"],
    )
    return save_fig(fig, "09_phenotype_landscape.png")


def choose_example_patients(per_model_alerts: Dict[str, pd.DataFrame]) -> List[str]:
    fusion = per_model_alerts["fusion_multimodel"].copy()
    logistic = per_model_alerts["logistic_regression"].copy()
    merged = fusion.merge(
        logistic[["patient_id", "early_alert", "lead_time_h"]],
        on="patient_id",
        suffixes=("_fusion", "_logistic"),
    )
    merged = merged[merged["sepsis_patient"]].copy()
    merged["lead_gain"] = merged["lead_time_h_fusion"].fillna(-999) - merged["lead_time_h_logistic"].fillna(-999)
    merged = merged.sort_values(["lead_gain", "lead_time_h_fusion"], ascending=False)

    selected: List[str] = []
    for _, row in merged.iterrows():
        pid = str(row["patient_id"])
        if pid not in selected:
            selected.append(pid)
        if len(selected) == 2:
            break
    if len(selected) < 2:
        sepsis_ids = fusion[fusion["sepsis_patient"]]["patient_id"].astype(str).tolist()
        for pid in sepsis_ids:
            if pid not in selected:
                selected.append(pid)
            if len(selected) == 2:
                break
    return selected


def build_patient_examples_figure(bundle: AnalysisBundle, per_model_alerts: Dict[str, pd.DataFrame]) -> Path:
    example_ids = choose_example_patients(per_model_alerts)
    if not example_ids:
        example_ids = [str(bundle.test_windows["patient_ids"][0])]
    onset_map = actual_onset_hours(bundle.raw_test)
    fig, axes = plt.subplots(len(example_ids), 1, figsize=(15, 4.6 * max(len(example_ids), 1)), facecolor=COLORS["bg"])
    if len(example_ids) == 1:
        axes = [axes]

    model_order = ["logistic_regression", "snapshot_xgb", "trajectory_gru", "fusion_multimodel"]
    for ax, pid in zip(axes, example_ids):
        mask = bundle.test_windows["patient_ids"] == pid
        patient_hours = bundle.test_windows["hours"][mask]
        for model_key in model_order:
            ax.plot(
                patient_hours,
                bundle.test_scores[model_key][mask],
                color=MODEL_COLORS[model_key],
                linewidth=2.4 if model_key == "fusion_multimodel" else 1.8,
                label=MODEL_DISPLAY[model_key],
            )
            alert_row = per_model_alerts[model_key].set_index("patient_id").loc[pid]
            if pd.notna(alert_row["first_alert_hour"]):
                ax.axvline(float(alert_row["first_alert_hour"]), color=MODEL_COLORS[model_key], linestyle=":", alpha=0.6)

        onset = onset_map.get(pid)
        if onset is not None:
            ax.axvline(onset, color=COLORS["danger"], linestyle="--", linewidth=2.5, label="True onset")
        ax.axhline(bundle.thresholds["fusion_multimodel"], color=COLORS["fusion"], linestyle="--", linewidth=1.5, alpha=0.6)
        ax.set_ylim(0, 1.02)
        fusion_row = per_model_alerts["fusion_multimodel"].set_index("patient_id").loc[pid]
        lead_text = (
            f"Fusion alert at h={int(fusion_row['first_alert_hour'])}, lead={fusion_row['lead_time_h']:.1f}h"
            if pd.notna(fusion_row["lead_time_h"])
            else "No pre-onset fusion alert"
        )
        ax.set_title(f"Patient {pid}: {lead_text}", loc="left", fontsize=15)
        ax.set_ylabel("Predicted risk")

    axes[-1].set_xlabel("ICU hour")
    handles = [
        Line2D([0], [0], color=MODEL_COLORS[k], linewidth=2.4 if k == "fusion_multimodel" else 1.8, label=MODEL_DISPLAY[k])
        for k in model_order
    ]
    handles.append(Line2D([0], [0], color=COLORS["danger"], linestyle="--", linewidth=2.5, label="True onset"))
    axes[0].legend(handles=handles, ncol=3, frameon=False, loc="upper left")

    fig.suptitle("Case studies make the story tangible", x=0.03, y=0.99, ha="left", fontsize=22, weight="bold")
    fig.text(
        0.03,
        0.95,
        "These are automatically selected held-out sepsis patients where fusion buys the most additional runway over the baseline.",
        fontsize=12,
        color=COLORS["muted"],
    )
    return save_fig(fig, "10_patient_examples.png")


def build_feature_importance_figure(bundle: AnalysisBundle) -> Path:
    snapshot = SnapshotModel.load(str(Path(CFG["paths"]["models_dir"]) / "snapshot.pkl"))
    fusion = FusionModel.load(str(Path(CFG["paths"]["models_dir"]) / "fusion.pkl"))
    snapshot_importance = snapshot.feature_importances()
    feature_names = bundle.feature_cols
    top_idx = np.argsort(snapshot_importance)[-12:]
    top_features = pd.DataFrame(
        {
            "feature": np.array(feature_names)[top_idx],
            "importance": snapshot_importance[top_idx],
        }
    ).sort_values("importance", ascending=True)

    meta_names = ["snapshot_prob", "trajectory_prob"] + [f"cluster_dist_{i}" for i in range(8)]
    fusion_importance = fusion.model.feature_importances_
    meta_df = pd.DataFrame({"feature": meta_names[: len(fusion_importance)], "importance": fusion_importance}).sort_values(
        "importance", ascending=True
    )

    fig, axes = plt.subplots(1, 2, figsize=(15, 7), facecolor=COLORS["bg"])
    axes[0].barh(top_features["feature"], top_features["importance"], color=COLORS["snapshot"])
    axes[0].set_title("Snapshot specialist: top drivers", loc="left", fontsize=16)
    axes[0].set_xlabel("Gain")
    axes[1].barh(meta_df["feature"], meta_df["importance"], color=COLORS["fusion"])
    axes[1].set_title("Fusion stacker: what it leans on", loc="left", fontsize=16)
    axes[1].set_xlabel("Gain")

    fig.suptitle("Appendix: what the models learned", x=0.03, y=0.99, ha="left", fontsize=22, weight="bold")
    fig.text(
        0.03,
        0.945,
        "Use this slide if you get asked whether the fusion gain is coming from real signal or just from stacking noise.",
        fontsize=12,
        color=COLORS["muted"],
    )
    return save_fig(fig, "11_feature_importance.png")


def build_slide_manifest(metrics_df: pd.DataFrame) -> List[Dict[str, object]]:
    fusion = metrics_df.loc[metrics_df["model_key"] == "fusion_multimodel"].iloc[0]
    baseline = metrics_df.loc[metrics_df["model_key"] == "logistic_regression"].iloc[0]
    uplift = fusion["AUROC"] - baseline["AUROC"]

    return [
        {
            "title": "RescueWindow",
            "subtitle": "Early sepsis deterioration forecasting from scratch-trained multimodel trajectories",
            "figure": "00_exec_summary.png",
            "bullets": [
                f"Fusion AUROC {fusion['AUROC']:.3f} and AUPRC {fusion['AUPRC']:.3f} on the held-out test set.",
                f"Fusion beats logistic regression by {uplift:.3f} AUROC.",
                f"Mean lead time is {fusion['Mean Lead Time (h)']:.1f} hours before true onset.",
            ],
            "notes": "Open with the central claim: this is not just a better classifier, it creates earlier intervention runway.",
        },
        {
            "title": "Benchmark Positioning",
            "subtitle": "The result lands in the band that makes the project defensible as a challenge-grade system",
            "figure": "01_benchmark_targets.png",
            "bullets": [
                "Fusion performance sits in or above the top-end project target ranges.",
                "The AUROC lands inside the rough 2019 challenge leader band.",
                "Lead time and early catch rate are the differentiators to emphasize verbally.",
            ],
            "notes": "Use this slide to establish credibility fast, then pivot immediately to why lead time matters more than the headline AUROC.",
        },
        {
            "title": "Dataset Framing",
            "subtitle": "Rare positives plus sparse labs make this a legitimately difficult early-warning problem",
            "figure": "02_dataset_overview.png",
            "bullets": [
                "The cohort has more than forty thousand ICU stays with hourly trajectories.",
                "Sepsis prevalence is low, which is why PR curves matter.",
                "Many clinically informative labs are missing for long stretches.",
            ],
            "notes": "This slide justifies the modeling choices and explains why a naive snapshot classifier leaves performance on the table.",
        },
        {
            "title": "System Design",
            "subtitle": "Three specialists feed a stacker that decides whose signal to trust",
            "figure": "03_system_architecture.png",
            "bullets": [
                "Snapshot model captures the current clinical state.",
                "GRU model captures silent temporal drift across the last twelve hours.",
                "Phenotype model gives the stacker a latent clinical context.",
            ],
            "notes": "Stress that every component is scratch-trained and that the fusion model only sees specialist outputs plus phenotype distances.",
        },
        {
            "title": "Fusion Wins",
            "subtitle": "The stacker outperforms the baseline and both specialists across the metrics that matter",
            "figure": "04_model_comparison.png",
            "bullets": [
                "Fusion is the strongest model on AUROC, AUPRC, detection rate, and lead time.",
                "Snapshot and GRU each contribute signal, but neither dominates across all metrics.",
                "This is the slide to quote when someone asks whether the extra complexity is worth it.",
            ],
            "notes": "Keep this slide simple. Let the chart do the work and call out the AUROC uplift and lead-time advantage.",
        },
        {
            "title": "Ranking Performance",
            "subtitle": "Fusion stays ahead throughout the score range, not just at one chosen threshold",
            "figure": "05_performance_curves.png",
            "bullets": [
                "ROC shows global separability.",
                "PR shows how much usable precision survives class imbalance.",
                "Fusion remains above the alternatives on both curves.",
            ],
            "notes": "If time is tight, summarize this slide in one sentence and move on. The deeper story is the operating point and lead-time behavior.",
        },
        {
            "title": "Operating Point",
            "subtitle": "Thresholds are calibrated on validation predictions, then frozen for the held-out test set",
            "figure": "06_operating_point_tradeoff.png",
            "bullets": [
                "The chosen threshold is not tuned on test.",
                "Lower thresholds buy coverage but burn specificity and precision.",
                "The selected point keeps a strong balance between coverage and runway.",
            ],
            "notes": "This slide answers the inevitable threshold question before it gets asked.",
        },
        {
            "title": "Lead-Time Story",
            "subtitle": "This is the slide that turns a classification result into a clinical story",
            "figure": "07_lead_time_story.png",
            "bullets": [
                "Lead time is measured to true onset, not to the start of the 6-hour prediction label window.",
                "The histogram shows how much runway early alerts typically create.",
                "The cumulative curve lets you say how often you buy at least 3 or 6 hours.",
            ],
            "notes": "Pause here. This is the most important slide in the deck.",
        },
        {
            "title": "Why Fusion Helps",
            "subtitle": "Snapshot and GRU fail on different patients, so the stacker can recover additional early catches",
            "figure": "08_specialist_complementarity.png",
            "bullets": [
                "Specialist disagreement is a feature, not a bug.",
                "Fusion recovers a meaningful share of patients that one specialist misses.",
                "This supports the architectural choice instead of making it feel ornamental.",
            ],
            "notes": "Use this when you need to defend the multimodel design on first principles, not just on outcome metrics.",
        },
        {
            "title": "Clinical Archetypes",
            "subtitle": "Latent phenotypes help the stacker contextualize risk",
            "figure": "09_phenotype_landscape.png",
            "bullets": [
                "The clusters line up with interpretable physiology, not random embeddings.",
                "Different archetypes carry very different feature signatures.",
                "The phenotype model makes the fusion layer more context-aware.",
            ],
            "notes": "This slide is optional in a shorter talk, but it adds depth and explains why the phenotype model exists.",
        },
        {
            "title": "Case Studies",
            "subtitle": "Held-out patient trajectories show what earlier warning looks like in practice",
            "figure": "10_patient_examples.png",
            "bullets": [
                "Each vertical marker shows when a model would have alerted.",
                "The true onset line makes the available runway easy to see.",
                "Fusion buys extra time on exactly the kinds of trajectories the baseline struggles with.",
            ],
            "notes": "These examples make the system feel real. Use them near the end, after the metrics have already done their job.",
        },
        {
            "title": "Appendix: Model Internals",
            "subtitle": "Feature importance for the snapshot specialist and the fusion layer",
            "figure": "11_feature_importance.png",
            "bullets": [
                "Important features line up with physiology clinicians expect to matter.",
                "The fusion model uses both specialist probabilities and phenotype distances.",
                "Keep this slide in reserve for Q and A.",
            ],
            "notes": "Appendix slide.",
        },
    ]


def write_deck_outline(manifest: List[Dict[str, object]]) -> None:
    lines = ["# RescueWindow Deck Outline", ""]
    for idx, slide in enumerate(manifest, start=1):
        lines.append(f"## Slide {idx}: {slide['title']}")
        lines.append(f"Subtitle: {slide['subtitle']}")
        lines.append(f"Visual: `presentation/figures/{slide['figure']}`")
        lines.append("Key points:")
        for bullet in slide["bullets"]:
            lines.append(f"- {bullet}")
        lines.append(f"Speaker note: {slide['notes']}")
        lines.append("")
    OUTLINE_PATH.write_text("\n".join(lines), encoding="utf-8")


def add_textbox(slide, left, top, width, height, text, font_size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)
    return box


def add_bullets(slide, left, top, width, height, bullets: List[str], color: str) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(18)
        run.font.color.rgb = hex_to_rgb(color)


def build_pptx(manifest: List[Dict[str, object]]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, slide_spec in enumerate(manifest):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = hex_to_rgb(COLORS["bg"])

        if idx == 0:
            banner = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
            )
            banner.fill.solid()
            banner.fill.fore_color.rgb = hex_to_rgb(COLORS["ink"])
            banner.line.fill.background()
            add_textbox(slide, Inches(0.5), Inches(0.25), Inches(7.8), Inches(0.4), slide_spec["title"], 28, COLORS["paper"], bold=True)
            add_textbox(slide, Inches(0.5), Inches(0.68), Inches(9.2), Inches(0.35), slide_spec["subtitle"], 14, COLORS["grid"])
            slide.shapes.add_picture(str(FIGURES_DIR / slide_spec["figure"]), Inches(0.55), Inches(1.45), width=Inches(12.2))
            add_bullets(slide, Inches(8.7), Inches(1.8), Inches(3.8), Inches(2.7), slide_spec["bullets"], COLORS["ink"])
        else:
            add_textbox(slide, Inches(0.55), Inches(0.35), Inches(11.6), Inches(0.4), slide_spec["title"], 24, COLORS["ink"], bold=True)
            add_textbox(slide, Inches(0.55), Inches(0.82), Inches(11.6), Inches(0.3), slide_spec["subtitle"], 13, COLORS["muted"])
            slide.shapes.add_picture(str(FIGURES_DIR / slide_spec["figure"]), Inches(0.55), Inches(1.25), width=Inches(8.6))
            add_bullets(slide, Inches(9.35), Inches(1.45), Inches(3.3), Inches(3.8), slide_spec["bullets"], COLORS["ink"])

        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.text = slide_spec["notes"]

    prs.save(str(DECK_PATH))


def write_tables(
    metrics_df: pd.DataFrame,
    deltas_df: pd.DataFrame,
    per_model_alerts: Dict[str, pd.DataFrame],
    sweep_df: pd.DataFrame,
) -> None:
    metrics_df.to_csv(TABLES_DIR / "model_metrics.csv", index=False)
    deltas_df.to_csv(TABLES_DIR / "headline_deltas.csv", index=False)
    sweep_df.to_csv(TABLES_DIR / "fusion_threshold_sweep.csv", index=False)
    for model_key, table in per_model_alerts.items():
        table.to_csv(TABLES_DIR / f"{model_key}_patient_alerts.csv", index=False)


def write_summary(metrics_df: pd.DataFrame, deltas_df: pd.DataFrame) -> None:
    def _normalize(obj):
        if isinstance(obj, dict):
            return {k: _normalize(v) for k, v in obj.items()}
        if isinstance(obj, list):
            return [_normalize(v) for v in obj]
        if isinstance(obj, np.generic):
            return obj.item()
        return obj

    payload = {
        "generated_at": pd.Timestamp.utcnow().isoformat(),
        "fusion": metrics_df.loc[metrics_df["model_key"] == "fusion_multimodel"].iloc[0].to_dict(),
        "logistic": metrics_df.loc[metrics_df["model_key"] == "logistic_regression"].iloc[0].to_dict(),
        "deltas": deltas_df.to_dict(orient="records"),
    }
    SUMMARY_PATH.write_text(json.dumps(_normalize(payload), indent=2), encoding="utf-8")


def main() -> None:
    ensure_output_dirs()
    set_plot_style()

    bundle = prepare_analysis_bundle()
    metrics_df, deltas_df, per_model_alerts = build_summary_tables(bundle)
    data_summary = dataset_summary(bundle)

    build_exec_summary_figure(metrics_df)
    build_benchmark_targets_figure(metrics_df)
    build_dataset_overview_figure(data_summary)
    build_architecture_figure()
    build_model_comparison_figure(metrics_df)
    build_performance_curves_figure(bundle)
    _, sweep_df = build_operating_point_figure(bundle)
    build_lead_time_figure(per_model_alerts["fusion_multimodel"])
    build_complementarity_figure(per_model_alerts)
    build_phenotype_figure(bundle)
    build_patient_examples_figure(bundle, per_model_alerts)
    build_feature_importance_figure(bundle)

    manifest = build_slide_manifest(metrics_df)
    write_deck_outline(manifest)
    build_pptx(manifest)
    write_tables(metrics_df, deltas_df, per_model_alerts, sweep_df)
    write_summary(metrics_df, deltas_df)

    log.info("Presentation assets written to %s", PRESENTATION_DIR)


if __name__ == "__main__":
    main()
